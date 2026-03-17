import os
import shutil
import httpx
from datetime import datetime
from pathlib import Path
from typing import List, Optional

import sys
from loguru import logger
from fastapi import UploadFile, HTTPException
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
import base64
import fitz  # PyMuPDF

import json
import uuid
from bson import ObjectId
from openai import AsyncOpenAI
from app.modules.knowledge.models import KnowledgeDocument, KnowledgeChatSession, KnowledgeChatMessage, KnowledgeNote
from app.modules.knowledge.schemas import ChatRequest, ChatResponse, ChatSessionResponse, ChatMessageSchema, NoteCreate, NoteResponse, PodcastResponse, InfographicRequest, InfographicResponse

from app.odm.user import UserDocument
from app.odm.system_settings import SystemSettings

# Initialize paths
UPLOAD_DIR = Path("app/static/uploads/knowledge")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
PODCAST_DIR = Path("app/static/podcasts")
PODCAST_DIR.mkdir(parents=True, exist_ok=True)




async def call_ai(system_prompt: str, user_prompt: str, service_name: str = "chat") -> str:
    """
    Helper to call AI API using dynamic settings.
    Supports OpenAI-compatible (LiteLLM) providers.
    """
    settings = await SystemSettings.get_settings()
    api_key = (settings.ai_api_key or os.getenv("LITELLM_API_KEY", "")).strip() or None
    base_url = (settings.ai_base_url or os.getenv("LITELLM_BASE_URL", "")).strip() or None
    
    model_name = settings.model_mappings.get(service_name, "gpt-4o")
    
    logger.info(f"DEBUG AI: Calling AI with model={model_name}, service={service_name}")
    logger.info(f"DEBUG AI: base_url={base_url}")
    if not api_key:
        logger.warning("DEBUG AI: api_key is missing")
        return "⚠️ AI API Key not configured. Please add it to your Admin Settings."

    if not base_url or not api_key:
        return "⚠️ AI API Key or Base URL not configured. Please check your admin settings."

    try:
        client = AsyncOpenAI(api_key=api_key, base_url=base_url)
        response = await client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.3
        )
        return response.choices[0].message.content
                
    except Exception as e:
        logger.error(f"AI Call failed: {e}")
        return f"⚠️ AI process failed: {str(e)}"

async def call_vision_ai(system_prompt: str, user_prompt: str, image_data_base64: str) -> str:
    """
    Helper to call AI API with image data (Vision/OCR).
    """
    settings = await SystemSettings.get_settings()
    api_key = (settings.ai_api_key or os.getenv("LITELLM_API_KEY", "")).strip() or None
    base_url = (settings.ai_base_url or os.getenv("LITELLM_BASE_URL", "")).strip() or None
    
    model_name = settings.model_mappings.get("vision", "Qwen/Qwen3-VL-235B")
    
    logger.info(f"DEBUG VISION: Calling Vision AI with model={model_name}")
    if not api_key:
        return "⚠️ AI API Key not configured."

    if not base_url or not api_key:
        return "⚠️ AI API Key or Base URL not configured for Vision."

    try:
        client = AsyncOpenAI(api_key=api_key, base_url=base_url)
        response = await client.chat.completions.create(
            model=model_name,
            messages=[
                {"role": "system", "content": system_prompt},
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": user_prompt},
                        {
                            "type": "image_url",
                            "image_url": {"url": f"data:image/jpeg;base64,{image_data_base64}"}
                        }
                    ]
                }
            ],
            max_tokens=2048
        )
        return response.choices[0].message.content
                
    except Exception as e:
        logger.error(f"Primary Vision Call failed ({model_name}): {e}")
        
        # Fallback to Google Gemini if available and primary fails
        google_key = os.getenv("GOOGLE_API_KEY")
        if google_key and ("401" in str(e) or "access denied" in str(e).lower() or "not found" in str(e).lower()):
            logger.info("OCR: Attempting fallback to Google Gemini Vision...")
            try:
                from langchain_google_genai import ChatGoogleGenerativeAI
                from langchain_core.messages import HumanMessage
                
                llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash", google_api_key=google_key)
                
                message = HumanMessage(
                    content=[
                        {"type": "text", "text": f"{system_prompt}\n\n{user_prompt}"},
                        {
                            "type": "image_url",
                            "image_url": f"data:image/jpeg;base64,{image_data_base64}",
                        },
                    ]
                )
                res = await llm.ainvoke([message])
                return res.content
            except Exception as ge:
                logger.error(f"Gemini fallback failed: {ge}")
                return f"⚠️ OCR failed. Primary: {str(e)}. Fallback: {str(ge)}"
        
        return f"⚠️ Vision process failed: {str(e)}"

# General purpose AI call (formerly call_gemini)
async def call_general_ai(system_prompt: str, user_prompt: str) -> str:
    return await call_ai(system_prompt, user_prompt, "chat")


class KnowledgeService:
    async def ingest_document(self, file: UploadFile, user: UserDocument) -> KnowledgeDocument:
        """Process uploaded PDF: Save -> Extract Text -> Store record"""
        try:
            timestamp = int(datetime.utcnow().timestamp())
            safe_filename = f"{timestamp}_{file.filename}"
            file_path = UPLOAD_DIR / safe_filename

            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            # Extract Text based on file extension
            logger.info(f"DEBUG: Starting ingestion for {file.filename}")
            
            from langchain_core.documents import Document as LcDocument
            docs = []
            
            if file.filename.lower().endswith(".pdf"):
                loader = PyMuPDFLoader(str(file_path))
                docs = loader.load()
                logger.info(f"DEBUG: Loaded {len(docs)} pages from PDF {file.filename}")
            elif file.filename.lower().endswith((".jpg", ".jpeg", ".png", ".webp")):
                image_text = await self.ocr_extract_image_text(file_path)
                if image_text:
                    docs = [LcDocument(page_content=image_text, metadata={"source": str(file_path)})]
                    logger.info(f"DEBUG: Extracted text from image {file.filename}")
                else:
                    docs = []
            elif file.filename.lower().endswith((".docx", ".doc")):
                import docx
                doc_obj = docx.Document(file_path)
                text_content = "\n".join([p.text for p in doc_obj.paragraphs if p.text.strip()])
                docs = [LcDocument(page_content=text_content, metadata={"source": str(file_path)})]
                logger.info(f"DEBUG: Extracted text from DOCX {file.filename}")
            elif file.filename.lower().endswith((".pptx", ".ppt")):
                import pptx
                prs = pptx.Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text_content = "\n".join(text_runs)
                docs = [LcDocument(page_content=text_content, metadata={"source": str(file_path)})]
                logger.info(f"DEBUG: Extracted text from PPTX {file.filename}")
            else:
                raise HTTPException(status_code=400, detail="Unsupported file format. Please upload PDF, DOCX, or PPTX.")

            text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000, chunk_overlap=200)
            splits = text_splitter.split_documents(docs)
            chunk_texts = [s.page_content.strip() for s in splits if s.page_content.strip()]
            logger.info(f"DEBUG: Total extracted text length: {sum(len(c) for c in chunk_texts)}")
            logger.info(f"DEBUG: Split into {len(chunk_texts)} non-empty chunks")

            if not chunk_texts:
                logger.warning(f"DEBUG: No text extracted from {file.filename}. Attempting OCR...")
                ocr_text = await self.ocr_extract_text(file_path)
                if ocr_text:
                    logger.info(f"DEBUG: OCR Extraction successful ({len(ocr_text)} chars)")
                    docs = [LcDocument(page_content=ocr_text, metadata={"source": str(file_path)})]
                    splits = text_splitter.split_documents(docs)
                    chunk_texts = [s.page_content.strip() for s in splits if s.page_content.strip()]
                    status = "processed"
                    error_message = None
                else:
                    logger.warning(f"DEBUG: OCR Extraction failed for {file.filename}")
                    status = "error"
                    # Pass the error from the AI if it failed due to config
                    error_message = "No text could be extracted. This is likely a scanned document and your Vision/OCR AI settings (API Key or Model) might be incorrect."
            else:
                status = "processed"
                error_message = None

            doc = await KnowledgeDocument.create(
                title=file.filename,
                filename=file.filename,
                file_size=file.size or 0,
                content_type=file.content_type or "application/pdf",
                storage_path=str(file_path),
                chunk_count=len(chunk_texts),
                chunks=chunk_texts,
                status=status,
                error_message=error_message,
                user_id=user.uuid,
            )
            logger.info(f"DEBUG: Created document {doc.id} with status {doc.status}")
            return doc
        except Exception as e:
            logger.error(f"Ingestion failed: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to process document: {str(e)}")

    async def ocr_extract_text(self, file_path: Path) -> Optional[str]:
        """Convert PDF pages to images and extract text using Vision AI (Qwen-VL/Gemini)"""
        if not file_path.suffix.lower() == ".pdf":
            return None
            
        try:
            logger.info(f"Attempting OCR for {file_path}")
            doc = fitz.open(str(file_path))
            full_text = []
            
            # Limit OCR to first 10 pages for performance and cost
            max_pages = min(len(doc), 10)
            logger.info(f"OCR: Processing {max_pages} pages for {file_path.name}")
            
            for i in range(max_pages):
                page = doc.load_page(i)
                # Render page to image (300 DPI for high quality OCR)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_bytes = pix.tobytes("jpeg")
                b64_img = base64.b64encode(img_bytes).decode('utf-8')
                
                system_prompt = "You are a high-precision OCR agent specialized in reading scanned documents and handwriting."
                user_prompt = f"Accurately extract and transcribe ALL text from this page. Maintain formatting where possible. This is page {i+1} of the document."
                
                logger.info(f"OCR: Calling Vision AI for page {i+1}...")
                page_text = await call_vision_ai(system_prompt, user_prompt, b64_img)
                if page_text and not page_text.startswith("⚠️"):
                    logger.info(f"OCR: Successfully extracted text for page {i+1} ({len(page_text)} chars)")
                    full_text.append(f"--- PAGE {i+1} ---\n{page_text}")
                else:
                    logger.error(f"OCR failed for page {i+1}: Result={page_text}")
                    
            doc.close()
            return "\n\n".join(full_text) if full_text else None
        except Exception as e:
            logger.error(f"OCR Pipeline failed: {e}")
            return None

    async def ocr_extract_image_text(self, file_path: Path) -> Optional[str]:
        """Directly extract text from an image file using Vision AI"""
        try:
            with open(file_path, "rb") as image_file:
                img_bytes = image_file.read()
                b64_img = base64.b64encode(img_bytes).decode('utf-8')
                
            system_prompt = "You are a high-precision OCR agent specialized in reading handwritten notes and scanned images."
            user_prompt = "Transcribe all text from this image as accurately as possible. If it's a handwritten note, maintain the structure and intent."
            
            text = await call_vision_ai(system_prompt, user_prompt, b64_img)
            if text and not text.startswith("⚠️"):
                return text
            return None
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return None

    async def chat(self, request: ChatRequest, user: UserDocument) -> ChatResponse:
        """RAG Chat: Load user document chunks -> Gemini REST API"""
        if not request.message.strip():
            return ChatResponse(answer="Please ask a question.", sources=[])

        query = {"status": "processed", "user_id": user.uuid}
            
        if request.document_ids:
            query["_id"] = {"$in": [ObjectId(id) for id in request.document_ids]}
            
        user_docs = await KnowledgeDocument.find(query)
        print(f"DEBUG: Found {len(user_docs)} processed documents for user {user.uuid}")

        sources = []
        context_text = ""

        if user_docs:
            sources = [d.filename for d in user_docs]
            all_chunks: List[str] = []
            for d in user_docs:
                if d.chunks:
                    print(f"DEBUG: Adding {len(d.chunks)} chunks from {d.filename}")
                    all_chunks.extend(d.chunks)
                else:
                    print(f"DEBUG: Document {d.filename} has NO chunks stored!")
            
            # Use a slightly smaller chunk for context to stay within token limits safely
            context_text = "\n\n---\n\n".join(all_chunks)[:12000]
            print(f"DEBUG: Final context_text length: {len(context_text)}")

        # 1. Load Session History
        session = None
        session_history_text = ""
        
        if request.session_id:
            session = await KnowledgeChatSession.find_by_id(request.session_id)
            if not session or session.user_id != user.uuid:
                raise HTTPException(status_code=404, detail="Chat session not found")
            
            # Build memory string from past messages (last 10 turns to save tokens)
            memory_msgs = []
            for msg in session.messages[-10:]:
                role_label = "User" if msg.role == "user" else "Assistant"
                memory_msgs.append(f"{role_label}: {msg.content}")
            
            if memory_msgs:
                session_history_text = "\n".join(memory_msgs)

        if context_text:
            system_prompt = (
                "You are an intelligent assistant. Use the document context provided below to answer precisely. "
                "If the answer is not in the context, say you don't know based on the documents.\n\n"
                f"DOCUMENT CONTEXT:\n{context_text}"
            )
        else:
            system_prompt = (
                "You are a helpful AI assistant. The user has not uploaded any documents yet. "
                "Explain that they can upload PDFs to chat with their documents, but otherwise you can assist with general queries."
            )
            
        if session_history_text:
            system_prompt += f"\n\nCONVERSATION HISTORY:\n{session_history_text}"

        answer = await call_general_ai(system_prompt, request.message)
        
        # 2. Persist to Session
        if not session:
            # Generate a very simple title from message
            new_title = request.message[:30] + "..." if len(request.message) > 30 else request.message
            session = await KnowledgeChatSession.create(
                title=new_title,
                user_id=user.uuid,
                messages=[]
            )
            
        session.messages.append(KnowledgeChatMessage(role="user", content=request.message))
        session.messages.append(KnowledgeChatMessage(role="assistant", content=answer, sources=sources))
        session.updated_at = datetime.utcnow()
        await session.save()

        return ChatResponse(answer=answer, sources=sources, session_id=str(session.id))

    async def list_chat_sessions(self, user: UserDocument) -> List[ChatSessionResponse]:
        query = {"user_id": user.uuid}
        sessions = await KnowledgeChatSession.find(query, sort=[("updated_at", -1)])
        return [
            ChatSessionResponse(
                id=str(s.id),
                title=s.title,
                updated_at=s.updated_at,
                messages=[] # Don't load full messages for the list view
            )
            for s in sessions
        ]
        
    async def get_chat_session(self, session_id: str, user: UserDocument) -> ChatSessionResponse:
        session = await KnowledgeChatSession.find_by_id(session_id)
        if not session or session.user_id != user.uuid:
            raise HTTPException(status_code=404, detail="Chat session not found")
            
        return ChatSessionResponse(
            id=str(session.id),
            title=session.title,
            updated_at=session.updated_at,
            messages=[
                ChatMessageSchema(
                    role=m.role,
                    content=m.content,
                    sources=m.sources or [],
                    created_at=m.created_at
                ) for m in session.messages
            ]
        )
        
    async def delete_chat_session(self, session_id: str, user: UserDocument):
        session = await KnowledgeChatSession.find_by_id(session_id)
        if not session or session.user_id != user.uuid:
            raise HTTPException(status_code=404, detail="Chat session not found")
            
        await session.delete()
        return True

    async def rename_chat_session(self, session_id: str, title: str, user: UserDocument) -> ChatSessionResponse:
        session = await KnowledgeChatSession.find_by_id(session_id)
        if not session or session.user_id != user.uuid:
            raise HTTPException(status_code=404, detail="Chat session not found")
            
        session.title = title
        session.updated_at = datetime.utcnow()
        await session.save()
        
        return await self.get_chat_session(session_id, user)

    async def list_documents(self, user: UserDocument) -> List[KnowledgeDocument]:
        query = {"user_id": user.uuid}
        return await KnowledgeDocument.find(query, sort=[("created_at", -1)])

    async def delete_document(self, doc_id: str, user: UserDocument):
        doc = await KnowledgeDocument.find_by_id(doc_id)
        if not doc or doc.user_id != user.uuid:
            raise HTTPException(status_code=404, detail="Document not found")
        
        if doc.storage_path and os.path.exists(doc.storage_path):
            try:
                os.remove(doc.storage_path)
            except:
                pass
        
        await doc.delete()
        return True

    # --- Notes ---
    async def create_note(self, data: NoteCreate, user: UserDocument) -> NoteResponse:
        note = await KnowledgeNote.create(
            title=data.title,
            content=data.content,
            user_id=user.uuid
        )
        return NoteResponse(
            id=str(note.id),
            title=note.title,
            content=note.content,
            created_at=note.created_at,
            updated_at=note.updated_at
        )

    async def list_notes(self, user: UserDocument) -> List[NoteResponse]:
        notes = await KnowledgeNote.find({"user_id": user.uuid}, sort=[("updated_at", -1)])
        return [
            NoteResponse(
                id=str(n.id),
                title=n.title,
                content=n.content,
                created_at=n.created_at,
                updated_at=n.updated_at
            ) for n in notes
        ]

    async def delete_note(self, note_id: str, user: UserDocument):
        note = await KnowledgeNote.find_by_id(note_id)
        if not note or note.user_id != user.uuid:
            raise HTTPException(status_code=404, detail="Note not found")
            
        await note.delete()

    async def generate_audio_overview(self, document_ids: List[str], user: UserDocument, session_id: Optional[str] = None) -> PodcastResponse:
        """
        Generate a conversational audio overview of the selected documents.
        """
        # 1. Fetch document content
        docs = []
        for doc_id in document_ids:
            doc = await KnowledgeDocument.find_by_id(doc_id)
            if doc and str(doc.user_id) == str(user.uuid):
                docs.append(doc)
        
        if not docs:
            raise HTTPException(status_code=404, detail="No valid documents found")

        combined_text = "\n\n".join([ "\n".join(d.chunks[:5])[:2000] for d in docs]) # Limit per doc for performance

        # 2. Generate Dialogue Script with Gemini
        system_prompt = (
            "You are a scriptwriter for a popular learning podcast like NotebookLM's 'Deep Dive'. "
            "Generate a natural, engaging dialogue between two people: Host A (the expert) and Host B (the curious learner). "
            "They are discussing the provided documents. Make it sound like a real conversation with interjections, "
            "uh-huhs, and summaries of key points. Keep it concise (about 10-15 lines total).\n\n"
            "Format your response as a strict JSON array of objects: [{\"speaker\": \"A\", \"text\": \"...\"}, ...]"
        )
        
        user_message = f"Documents Content:\n{combined_text}"
        
        script_raw = await call_general_ai(system_prompt, user_message)
        
        try:
            # Clean possible markdown formatting from Gemini
            clean_script = script_raw.strip()
            if clean_script.startswith('```'):
                import re
                clean_script = re.sub(r'^```[a-z]*\n?', '', clean_script)
                clean_script = re.sub(r'\n?```$', '', clean_script)
            
            script = json.loads(clean_script)
        except Exception as e:
            logger.error(f"Failed to parse script JSON: {e}. Raw: {script_raw}")
            # Fallback to a simple script if generation fails
            script = [{"speaker": "A", "text": "I've analyzed your documents. They cover some very interesting topics!"}, 
                      {"speaker": "B", "text": "That's great! Can you tell me more about the key takeaways?"}]

        # 3. Generate Audio
        api_key_openai = os.getenv("PODCAST_API_KEY")
        api_key_elevenlabs = os.getenv("ELEVENLABS_API_KEY")
        
        if not api_key_openai and not api_key_elevenlabs:
            raise HTTPException(status_code=500, detail="Neither PODCAST_API_KEY nor ELEVENLABS_API_KEY is configured")
            
        combined_audio = b""
        filename = f"podcast_{user.uuid}_{uuid.uuid4().hex}.mp3"
        user_podcast_dir = PODCAST_DIR / str(user.uuid)
        user_podcast_dir.mkdir(parents=True, exist_ok=True)
        file_path = user_podcast_dir / filename

        try:
            if api_key_elevenlabs:
                # ElevenLabs Implementation
                async with httpx.AsyncClient() as client_http:
                    for line in script:
                        speaker = line.get("speaker", "A")
                        text = line.get("text", "")
                        # Verified IDs: Adam (Expert) and Rachel (Learner)
                        voice_id = "pNInz6obpgDQGcFmaJgB" if speaker == "A" else "21m00Tcm4TlvDq8ikWAM"
                        
                        url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"
                        headers = {
                            "Accept": "audio/mpeg",
                            "Content-Type": "application/json",
                            "xi-api-key": api_key_elevenlabs
                        }
                        data = {
                            "text": text,
                            "model_id": "eleven_flash_v2_5",
                            "voice_settings": {
                                "stability": 0.5,
                                "similarity_boost": 0.5
                            }
                        }
                        
                        response = await client_http.post(url, json=data, headers=headers)
                        if response.status_code != 200:
                            logger.error(f"ElevenLabs failed: {response.status_code} - {response.text}")
                            raise Exception(f"ElevenLabs API error: {response.status_code}")
                        
                        combined_audio += response.content
            else:
                # Fallback to OpenAI
                client = AsyncOpenAI(api_key=api_key_openai)
                for line in script:
                    speaker = line.get("speaker", "A")
                    text = line.get("text", "")
                    voice = "alloy" if speaker == "A" else "echo"
                    
                    response = await client.audio.speech.create(
                        model="tts-1",
                        voice=voice,
                        input=text
                    )
                    combined_audio += response.content

            with open(file_path, "wb") as f:
                f.write(combined_audio)
                
        except Exception as e:
            logger.error(f"Audio generation failed: {e}")
            raise HTTPException(status_code=500, detail=f"Failed to generate audio: {str(e)}")

        public_url = f"/static/podcasts/{user.uuid}/{filename}"
        
        # 4. Save to Chat Session
        session = None
        if session_id:
            session = await KnowledgeChatSession.find_by_id(session_id)
            if session and session.user_id != user.uuid:
                session = None # Invalid or not belonging to user

        if not session:
            session = await KnowledgeChatSession.create(
                title="Audio Podcast Overview",
                user_id=user.uuid,
                messages=[]
            )
            
        # Format the response exactly as the frontend would render it
        audio_block = f"[🎧 Listen to Podcast]({public_url})"
        transcript_text = "\n\n***\n\n### 🎙️ Discussion Transcript\n\n"
        for line in script:
            speaker_name = "Host" if line.get("speaker", "A") == "A" else "Guest"
            icon = "🎙️" if line.get("speaker", "A") == "A" else "🎧"
            clean_text = (line.get("text", "")).replace("\\n", " ")
            transcript_text += f"> **{icon} {speaker_name}:** {clean_text}\n>\n"
            
        full_content = f"{audio_block}{transcript_text}"
        
        from app.modules.knowledge.models import KnowledgeChatMessage
        session.messages.append(KnowledgeChatMessage(role="user", content="Generate an audio podcast discussing the selected documents."))
        
        source_filenames = [d.filename for d in docs]
        session.messages.append(KnowledgeChatMessage(role="assistant", content=full_content, sources=source_filenames))
        
        session.updated_at = datetime.utcnow()
        await session.save()

        return PodcastResponse(audio_url=public_url, script=script, session_id=str(session.id))

    async def generate_infographic(self, request: InfographicRequest, user: UserDocument) -> InfographicResponse:
        """
        Generate a data-driven infographic summary.
        1. Summarize documents into structured JSON with Gemini.
        2. Fetch a relevant theme image from Unsplash.
        3. Save as a Note.
        """
        # 1. Fetch document content
        docs = []
        for doc_id in request.document_ids:
            doc = await KnowledgeDocument.find_by_id(doc_id)
            if doc and str(doc.user_id) == str(user.uuid):
                docs.append(doc)
        
        if not docs:
            raise HTTPException(status_code=404, detail="No valid documents found")

        combined_text = "\n\n".join([ "\n".join(d.chunks[:10])[:4000] for d in docs]) 

        # 2. Generate Structured Summary with Gemini
        system_prompt = (
            "You are an expert data visualizer and infographic designer. "
            "Analyze the provided text and transform it into a premium, data-driven infographic structure.\n"
            "CRITICAL: YOUR RESPONSE MUST BE STOIC, RAW JSON ONLY. NO MARKDOWN CODE BLOCKS. NO INTRO OR OUTRO TEXT. JUST THE JSON OBJECT.\n"
            "Format your response as a strict JSON object with these EXACT keys:\n"
            "{\n"
            "  \"main_title\": \"(A short, punchy title)\",\n"
            "  \"intro\": \"(A 1-2 sentence compelling summary)\",\n"
            "  \"stats\": [{\"label\": \"(Short label)\", \"value\": \"(Number/Metric)\", \"unit\": \"(e.g. %, stars)\"}],\n"
            "  \"sections\": [\n"
            "    {\n"
            "      \"title\": \"(Section title)\", \n"
            "      \"content\": \"(Detailed finding)\",\n"
            "      \"icon_hint\": \"(Suggest a Lucide icon: zap, shield, database, activity, bar-chart, etc.)\"\n"
            "    }\n"
            "  ],\n"
            "  \"unsplash_query\": \"(A descriptive query for a clean, minimalist professional background image)\"\n"
            "}"
        )
        
        user_message = f"Documents Content:\n{combined_text}"
        summary_raw = await call_general_ai(system_prompt, user_message)
        
        try:
            clean_summary = summary_raw.strip()
            # Remove potential markdown code block wrappers
            if "```json" in clean_summary:
                clean_summary = clean_summary.split("```json")[1].split("```")[0]
            elif "```" in clean_summary:
                clean_summary = clean_summary.split("```")[1].split("```")[0]
            
            # Further clean any potential non-JSON text at start/end
            start_idx = clean_summary.find("{")
            end_idx = clean_summary.rfind("}")
            if start_idx != -1 and end_idx != -1:
                clean_summary = clean_summary[start_idx:end_idx+1]
                
            summary_data = json.loads(clean_summary)
        except Exception as e:
            logger.error(f"Failed to parse infographic JSON: {e}. Raw: {summary_raw}")
            summary_data = {
                "main_title": "Document Summary",
                "intro": "Analysed document content for key insights.",
                "stats": [{"label": "Documents", "value": str(len(docs)), "unit": "files"}],
                "sections": [{"title": "Overview", "content": "The documents provide a comprehensive look into the subject matter.", "icon_hint": "file-text"}],
                "unsplash_query": "abstract gradient minimalist background"
            }

        # 3. Fetch Image from Unsplash
        access_key = os.getenv("UNSPLASH_ACCESS_KEY")
        image_url = "https://images.unsplash.com/photo-1543286535-ea1967280e5b?auto=format&fit=crop&q=80&w=1080" # Fallback
        
        if access_key:
            keyword = summary_data.get("unsplash_query", "abstract background minimalist")
            try:
                async with httpx.AsyncClient() as client:
                    response = await client.get(
                        "https://api.unsplash.com/search/photos",
                        params={"query": keyword, "per_page": 1, "orientation": "landscape"},
                        headers={"Authorization": f"Client-ID {access_key}"}
                    )
                    if response.status_code == 200:
                        data = response.json()
                        if data["results"]:
                            image_url = data["results"][0]["urls"]["regular"]
            except Exception as e:
                logger.error(f"Unsplash API failed: {e}")

        # 4. Save to Note
        note_content = f"# 📊 Visual Infographic: {summary_data.get('main_title', 'Summary')}\n\n"
        note_content += f"![Background]({image_url})\n\n"
        note_content += f"> {summary_data.get('intro', '')}\n\n"
        note_content += "### 📈 Key Metrics\n"
        for stat in summary_data.get("stats", []):
            note_content += f"- **{stat['label']}:** {stat['value']} {stat.get('unit', '')}\n"
        note_content += "\n### 💡 Main Insights\n"
        for section in summary_data.get("sections", []):
            note_content += f"#### {section.get('title', 'Insight')}\n{section.get('content', '')}\n\n"

        await KnowledgeNote.create(
            title=f"Infographic: {summary_data.get('main_title', 'Summary')[:30]}",
            content=note_content,
            user_id=user.uuid
        )

        # 5. Connect to Session
        session = None
        if request.session_id:
            session = await KnowledgeChatSession.find_by_id(request.session_id)
        
        if not session:
            session = await KnowledgeChatSession.create(
                title=f"Infographic - {summary_data.get('main_title', 'Summary')[:20]}",
                user_id=user.uuid,
                messages=[]
            )

        infographic_msg = f"Generated an infographic for the selected documents.\n\n{note_content}"
        session.messages.append(KnowledgeChatMessage(role="user", content="Generate an infographic summary."))
        session.messages.append(KnowledgeChatMessage(role="assistant", content=infographic_msg, sources=[d.filename for d in docs]))
        session.updated_at = datetime.utcnow()
        await session.save()

        return InfographicResponse(
            image_url=image_url,
            summary_data=summary_data,
            session_id=str(session.id)
        )

    async def transcribe_audio(self, audio_file: UploadFile, user: UserDocument) -> str:
        """
        Transcribe audio using the local faster-whisper model.
        Converts to WAV via ffmpeg (preserving original quality) before transcribing.
        """
        import tempfile
        import subprocess

        try:
            content = await audio_file.read()
            logger.info(f"Received audio for transcription: size={len(content)} bytes, type={audio_file.content_type}")

            # Normalise MIME type
            raw_mime = (audio_file.content_type or "audio/webm").split(";")[0].strip()
            
            ext_map = {
                "audio/webm": "webm", "audio/ogg": "ogg",
                "audio/wav": "wav", "audio/wave": "wav",
                "audio/mpeg": "mp3", "audio/mp3": "mp3",
                "audio/mp4": "mp4", "audio/aac": "aac",
                "audio/flac": "flac",
            }
            src_ext = ext_map.get(raw_mime, "webm")
            
            with tempfile.NamedTemporaryFile(suffix=f".{src_ext}", delete=False) as tmp_in:
                tmp_in.write(content)
                tmp_in_path = tmp_in.name

            # DEBUG: Save a permanent copy to check if the audio is completely silent
            try:
                import shutil
                # Ensure directory exists
                os.makedirs("app/static/uploads/knowledge", exist_ok=True)
                debug_dest = f"app/static/uploads/knowledge/last_debug_audio.{src_ext}"
                shutil.copy2(tmp_in_path, debug_dest)
                logger.info(f"Saved debug audio to {debug_dest}. Size: {len(content)} bytes")
            except Exception as e:
                logger.error(f"Failed to save debug audio: {e}")

            # Prepare WAV version for Custom Whisper API (supports WAV/MP3 only)
            tmp_wav_path = tmp_in_path + ".wav"
            logger.info(f"Target WAV path: {tmp_wav_path}")
            try:
                # Convert to WAV (16kHz, mono) for optimal ASR compatibility
                logger.info(f"Running FFmpeg: {tmp_in_path} -> {tmp_wav_path}")
                proc = subprocess.run(
                    ["ffmpeg", "-y", "-i", tmp_in_path, 
                     "-vn", 
                     "-ar", "16000", 
                     "-ac", "1", 
                     "-codec:a", "pcm_s16le", 
                     tmp_wav_path],
                    capture_output=True, text=True, timeout=30
                )
                logger.info(f"FFmpeg exit code: {proc.returncode}")
                if proc.returncode != 0:
                    logger.error(f"FFmpeg conversion for API failed: {proc.stderr}")
                    # Fallback to original file if conversion fails, though it might fail at API
                    tmp_wav_path = tmp_in_path
                else:
                    logger.info(f"Successfully converted audio to WAV for Whisper API. Size: {os.path.getsize(tmp_wav_path)} bytes")
            except Exception as e:
                logger.error(f"FFmpeg conversion error: {e}")
                tmp_wav_path = tmp_in_path

            # Custom Whisper STT Integration
            whisper_api_base = os.getenv("WHISPER_API_BASE")
            whisper_api_key = os.getenv("WHISPER_API_KEY")
            
            logger.info(f"Whisper API Base: {whisper_api_base}")
            if whisper_api_base and whisper_api_key:
                logger.info(f"Using Custom Whisper STT at {whisper_api_base}...")
                try:
                    async with httpx.AsyncClient() as client_http:
                        # Normalize base URL (ensure no trailing slash and append /audio/transcriptions)
                        base_url = whisper_api_base.rstrip("/")
                        url = f"{base_url}/audio/transcriptions"
                        logger.info(f"Requesting Whisper API: {url}")
                        
                        headers = {
                            "Authorization": f"Bearer {whisper_api_key}"
                        }
                        # Use the WAV file for the API call
                        api_src_path = tmp_wav_path
                        api_mime = "audio/wav" if tmp_wav_path.endswith(".wav") else raw_mime
                        
                        logger.info(f"Opening file: {api_src_path} as {api_mime}")
                        with open(api_src_path, "rb") as audio_f:
                            files = {
                                "file": (os.path.basename(api_src_path), audio_f, api_mime)
                            }
                            data = {
                                "model": "whisper-1", # Standard model name
                                "language": "en"
                            }
                            
                            logger.info("Sending POST request to Whisper API...")
                            response = await client_http.post(url, headers=headers, files=files, data=data, timeout=60)
                        
                        logger.info(f"Whisper API Status: {response.status_code}")
                        if response.status_code == 200:
                            res_json = response.json()
                            logger.info(f"Whisper API Response Type: {type(res_json)}")
                            logger.info(f"Whisper API Response Content: {res_json}")
                            
                            if isinstance(res_json, dict):
                                text = res_json.get("text", "").strip()
                            elif isinstance(res_json, list) and len(res_json) > 0:
                                # Try to join texts if it's a list of segments
                                text_parts = []
                                for item in res_json:
                                    if isinstance(item, dict):
                                        # Check common fields
                                        t = item.get("text") or item.get("metadata", {}).get("text")
                                        if t: text_parts.append(t)
                                text = " ".join(text_parts).strip()
                            else:
                                text = str(res_json).strip()

                            logger.info(f"Whisper STT Result: '{text}'")
                            
                            # Apply hallucination filter here too if we return early
                            clean_text = text.strip()
                            lower_text = clean_text.lower().replace(".", "").replace(",", "").replace("!", "").replace("?", "").strip()
                            hallucination_exact_matches = [
                                "you", "thank you", "thanks", "thank you for watching", 
                                "watching", "subtitle by", "you i", "yo", "peace", "bye", "go"
                            ]
                            if not lower_text or lower_text in hallucination_exact_matches:
                                logger.warning(f"Whisper API Hallucination or silence detected: '{clean_text}'")
                                return ""

                            # Clean up temporary WAV if it was created
                            if tmp_wav_path != tmp_in_path:
                                try: os.remove(tmp_wav_path)
                                except: pass
                            return clean_text
                        else:
                            logger.error(f"Whisper STT failed: {response.status_code} - {response.text}")
                            # Fallback to local whisper below
                except Exception as e:
                    logger.error(f"Whisper STT error, falling back to local: {e}")
                finally:
                    # Ensure tmp_wav_path is removed if it wasn't already
                    if tmp_wav_path != tmp_in_path and os.path.exists(tmp_wav_path):
                        try: os.remove(tmp_wav_path)
                        except: pass



            # Local Whisper Fallback
            tmp_out_path = tmp_in_path + "_out.wav"
            try:
                # Convert to WAV (16kHz, mono) for Whisper
                proc = subprocess.run(
                    ["ffmpeg", "-y", "-i", tmp_in_path, 
                     "-vn", 
                     "-ar", "16000", 
                     "-ac", "1", 
                     "-codec:a", "pcm_s16le", 
                     tmp_out_path],
                    capture_output=True, text=True, timeout=30
                )
                logger.info(f"ffmpeg (WAV) exit={proc.returncode}")
                if proc.returncode != 0:
                    logger.error(f"ffmpeg error: {proc.stderr}")
                    raise HTTPException(status_code=500, detail="Audio conversion failed")
                
                # Initialize local Whisper model (lazily inside the function for safety on module load)
                # It caches the model in memory via the backend singleton pattern, or just loads it here
                # if not previously cached. For a simple local approach, we'll try to keep it globally.
                global _whisper_model
                if '_whisper_model' not in globals() or _whisper_model is None:
                    from faster_whisper import WhisperModel
                    logger.info("Initializing faster-whisper model (base, cpu, int8)...")
                    _whisper_model = WhisperModel("base", device="cpu", compute_type="int8")
                    logger.info("faster-whisper model initialized.")
                
                # Transcribe
                logger.info("Starting local Whisper transcription...")
                segments, info = _whisper_model.transcribe(
                    tmp_out_path, 
                    beam_size=5, 
                    language="en", 
                    condition_on_previous_text=False,
                    vad_filter=True,
                    vad_parameters=dict(min_silence_duration_ms=2000, speech_pad_ms=400)
                )
                
                text_segments = []
                for segment in segments:
                    logger.info(f"Segment [{segment.start:.2f}s -> {segment.end:.2f}s]: {segment.text}")
                    text_segments.append(segment.text)
                    
                text = " ".join(text_segments).strip()

            finally:
                for p in [tmp_in_path, tmp_out_path]:
                    try: os.remove(p)
                    except Exception: pass

            logger.info(f"Local ASR Raw Result (len={len(text)}): '{text}'")
            
            clean_text = text.strip()
            # Filter out extremely common Whisper hallucinations on silence/noise
            lower_text = clean_text.lower().replace(".", "").replace(",", "").replace("!", "").replace("?", "").strip()
            
            hallucination_exact_matches = [
                "you", "thank you", "thanks", "thank you for watching", 
                "watching", "subtitle by", "you i", "yo", "peace", "bye", "go"
            ]
            
            if not lower_text or lower_text in hallucination_exact_matches:
                logger.warning(f"ASR Hallucination or silence detected: '{clean_text}'")
                return ""
            
            logger.info(f"Final Transcribed Text: '{clean_text}'")
            return clean_text

        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Transcription error: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Transcription error: {str(e)}")


knowledge_service = KnowledgeService()
